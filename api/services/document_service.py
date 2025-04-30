from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt
import os
from langchain_core.prompts import PromptTemplate
from models.llm_clients import LlmUtils

llm = LlmUtils.llm

TEXT_BLOCK_HEIGHT = Inches(0.8)
IMAGE_HEIGHT = Inches(2)
IMAGE_SPACING = Inches(0.2)
MAX_CONTENT_HEIGHT = Inches(7.0)
LOCAL_IMG_DIR="../assets/generated_images/2025-04-21-11-57-47"

class PptGenerator:


    def clean_request(request:str):
        final_req = request.replace("python","").replace("`", "")
        return final_req

    def generate_ppt_request(context: str):
        "generate a sample request for creating a ppt file using pptx library"

        prompt = """You are a smart assistant than can create payload samples for creating pptx files.
        The content of the sample must be extracted from this context:

        ### Context ###
        {context}
        Your job is to generate a sample request following the example below.
        

        The title field is the title of the document.
        The pages field contain a list of the different slides and their content.
        The TextItem contains two properties: type and content
            the type is a string and can have the following values: "title", "subtitle", "paragraph" 
            the content a string and represents the content of the text item
        
        The DocumentContent object is a list of TextItems.
        
        Please generate a sample request based on the context above. feel free to decide how many slides, the titles, paragraphs and content
        of each slide.

        ### Example of a request ###
        sample_request = DocumentRequest(
            title="Demo Document",
            pages=[
                DocumentContent(
                    text_items=[
                        TextItem(type="header", content="Welcome"),
                        TextItem(type="paragraph", content="This is a sample slide."),
                    ],
                ),

                DocumentContent(
                    text_items=[
                        TextItem(type="subheader", content="Second Slide"),
                        TextItem(type="paragraph", content="Lorem ipsum dolor sit amet, consectetur adipiscing elit. Vivamus luctus urna sed urna ultricies ac tempor dui sagittis. In condimentum facilisis porta.\nFusce sed felis eget velit aliquet faucibus. Praesent ac massa at ligula laoreet iaculis."),
                    ],  
                )
            ]
        )

        Return just the sample payload. Nothing else.
        """

        prompt_template = PromptTemplate.from_template(prompt)
        chain  = prompt_template | llm
        sample_request = chain.invoke({"context": context})
        return sample_request.content

    def generate_ppt(doc_data, output_path: str):
        prs = Presentation()

        for slide_data in doc_data.pages:
            slide = prs.slides.add_slide(prs.slide_layouts[5])

            # Customizing text dimensions
            num_blocks = len(slide_data.text_items or [])
            text_height_estimate = num_blocks * TEXT_BLOCK_HEIGHT
            textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8.5), text_height_estimate)
            tf = textbox.text_frame
            tf.word_wrap = True
            tf.auto_size = None#MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            for i, element in enumerate(slide_data.text_items or []):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.text = element.content
                para.space_after = Pt(10)
                font = para.font

                if element.type == "title":
                    font.size = Pt(32)
                    font.bold = True
                elif element.type == "subtitle":
                    font.size = Pt(24)
                    font.italic = True
                elif element.type == "paragraph":
                    font.size = Pt(16)
                else:
                    font.size = Pt(14)

            # Customizing images
            image_top = Inches(0.5) + text_height_estimate + Inches(0.5)

            for img_name in slide_data.images or []:
                image_path = os.path.join(LOCAL_IMG_DIR, img_name)
                if os.path.exists(image_path):
                    # Check if this image fits on the current slide
                    if image_top + IMAGE_HEIGHT > MAX_CONTENT_HEIGHT:
                        # Create a new slide and reset
                        slide = prs.slides.add_slide(prs.slide_layouts[5])
                        image_top = Inches(0.5)

                    try:
                        slide.shapes.add_picture(image_path, Inches(0.5), image_top, height=IMAGE_HEIGHT)
                        image_top += IMAGE_HEIGHT + IMAGE_SPACING
                    except Exception as e:
                        print(f"❌ Error adding new image: {e}")

        prs.save(output_path)

        